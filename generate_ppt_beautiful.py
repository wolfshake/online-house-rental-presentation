import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Please install python-pptx first: pip install python-pptx")
    exit(1)

# ════════════════════ PREMIUM DESIGN CONSTANTS ════════════════════
# Corporate Slate Palette
COLOR_PRIMARY_DARK  = RGBColor(15, 23, 42)     # #0F172A (Deep Slate Blue)
COLOR_LIGHT_BG      = RGBColor(248, 250, 252)  # #F8FAFC (Ultra-Light Cool Grey)
COLOR_ACCENT        = RGBColor(232, 114, 42)   # #E8722A (Warm Sunset Coral)
COLOR_TEXT_MAIN     = RGBColor(30, 41, 59)     # #1E293B (Charcoal Dark)
COLOR_TEXT_MUTED    = RGBColor(100, 116, 139)  # #64748B (Slate Muted)
COLOR_CARD_BORDER   = RGBColor(226, 232, 240)  # #E2E8F0 (Soft Edge Border)
COLOR_WHITE         = RGBColor(255, 255, 255)  # #FFFFFF (Crisp White)

# Secondary Module Accents
COLOR_TEAL_PANEL    = RGBColor(13, 148, 136)   # #0D9488 (Teal Muted)
COLOR_AMBER_PANEL   = RGBColor(180, 83, 9)     # #B45309 (Warm Amber)

FONT_FAMILY = 'Arial'

# ════════════════════ CORE CONFIGURATION ════════════════════
prs = Presentation()
prs.slide_width = Inches(13.33)  # Professional Widescreen 16:9 
prs.slide_height = Inches(7.5)

# ════════════════════ DESIGN HELPER FUNCTIONS ════════════════════
def apply_solid_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def draw_header_footer(slide, title_text, current_slide, total_slides=14, is_dark_layout=False):
    # Slide Title Text Frame
    if title_text:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.0), Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE if is_dark_layout else COLOR_PRIMARY_DARK

        # Clean Accent Underline Tracker
        accent_bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(1.2), Inches(0.04))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLOR_ACCENT
        accent_bar.line.fill.background()

    # Footer Tracker Line
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.4))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
    p_f = tf_f.paragraphs[0]
    p_f.text = f"Online House Rental System  |  Slide {current_slide} of {total_slides}"
    p_f.font.name = FONT_FAMILY
    p_f.font.size = Pt(9.5)
    p_f.font.color.rgb = RGBColor(148, 163, 184) if is_dark_layout else COLOR_TEXT_MUTED

def draw_bullet_point(slide, step_number, message_text, layout_top):
    # Minimalist Rounded Accent Circle
    badge = slide.shapes.add_shape(9, Inches(0.6), Inches(layout_top), Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_ACCENT
    badge.line.fill.background()
    tf_b = badge.text_frame
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.text = str(step_number)
    p_b.font.name = FONT_FAMILY
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_WHITE
    p_b.alignment = PP_ALIGN.CENTER

    # Paragraph Detail Text Box
    text_box = slide.shapes.add_textbox(Inches(1.15), Inches(layout_top - 0.05), Inches(11.5), Inches(0.6))
    tf_t = text_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = message_text
    p_t.font.name = FONT_FAMILY
    p_t.font.size = Pt(14)
    p_t.font.color.rgb = COLOR_TEXT_MAIN

def draw_container_card(slide, x, y, w, h, card_heading, bullet_items):
    # Geometric Structural Background Box
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_CARD_BORDER

    # Content Text Container
    text_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(w - 0.6), Inches(h - 0.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    # Header Card Text
    p_h = tf.paragraphs[0]
    p_h.text = card_heading
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(16)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_PRIMARY_DARK
    p_h.space_after = Pt(14)

    # Bullet Points Populating
    for text in bullet_items:
        p_li = tf.add_paragraph()
        p_li.text = "• " + text
        p_li.font.name = FONT_FAMILY
        p_li.font.size = Pt(12)
        p_li.font.color.rgb = COLOR_TEXT_MUTED
        p_li.space_after = Pt(8)

def draw_asymmetric_sidebar(slide, panel_bg_color, side_title, module_blocks, slide_index):
    # Deep Colored Column Block
    sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(3.2), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = panel_bg_color
    sidebar.line.fill.background()

    # Dynamic Capital Letter Label
    sb_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.2), Inches(2.4), Inches(3.0))
    tf_sb = sb_box.text_frame
    tf_sb.word_wrap = True
    p_sb = tf_sb.paragraphs[0]
    p_sb.text = side_title
    p_sb.font.name = FONT_FAMILY
    p_sb.font.size = Pt(28)
    p_sb.font.bold = True
    p_sb.font.color.rgb = COLOR_WHITE

    draw_header_footer(slide, None, slide_index, is_dark_layout=False)

    # Feature List Grid Placement
    for i, item in enumerate(module_blocks):
        y_offset = 1.6 + (i * 2.2)
        
        # Design Accent Bullet Node
        node = slide.shapes.add_shape(9, Inches(3.8), Inches(y_offset + 0.08), Inches(0.18), Inches(0.18))
        node.fill.solid()
        node.fill.fore_color.rgb = COLOR_ACCENT
        node.line.fill.background()

        # Item Text Wrapper
        item_box = slide.shapes.add_textbox(Inches(4.2), Inches(y_offset), Inches(8.5), Inches(1.8))
        tf_i = item_box.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0

        p_title = tf_i.paragraphs[0]
        p_title.text = item['title']
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_DARK
        p_title.space_after = Pt(4)

        p_desc = tf_i.add_paragraph()
        p_desc.text = item['desc']
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(12.5)
        p_desc.font.color.rgb = COLOR_TEXT_MAIN

def build_section_divider(slide, heavy_title, sub_headline, slide_index):
    apply_solid_background(slide, COLOR_PRIMARY_DARK)

    # Small Context Tag
    tag_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.33), Inches(0.4))
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = sub_headline.upper()
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.name = FONT_FAMILY
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_ACCENT

    # Bold Display Header Text
    h_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.9), Inches(11.33), Inches(1.2))
    tf_h = h_box.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = heavy_title
    p_h.alignment = PP_ALIGN.CENTER
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(44)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_WHITE

    # Minimalist Spacer Bar
    spacer_bar = slide.shapes.add_shape(1, Inches(5.66), Inches(4.4), Inches(2.0), Inches(0.04))
    spacer_bar.fill.solid()
    spacer_bar.fill.fore_color.rgb = COLOR_ACCENT
    spacer_bar.line.fill.background()

    draw_header_footer(slide, None, slide_index, is_dark_layout=True)

def build_relational_table(slide, x, y, w, h, table_header_lbl, headers, records):
    # Distinct Table Header Dark Label Cap
    cap = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.38))
    cap.fill.solid()
    cap.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    cap.line.fill.background()

    tf_c = cap.text_frame
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = f"  {table_header_lbl.upper()}"
    p_c.font.name = FONT_FAMILY
    p_c.font.size = Pt(11)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_WHITE

    # Generate PPT Grid Matrix
    row_count = len(records) + 1
    col_count = len(headers)
    grid_shape = slide.shapes.add_table(row_count, col_count, Inches(x), Inches(y + 0.38), Inches(w), Inches(h - 0.38))
    grid = grid_shape.table

    # Standard Schema Width Constraints
    proportional_widths = [0.22, 0.22, 0.56]
    for idx in range(col_count):
        grid.columns[idx].width = Inches(w * proportional_widths[idx])

    # Paint Schema Headers
    for c_idx, title in enumerate(headers):
        cell = grid.cell(0, c_idx)
        cell.text = title
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_FAMILY
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

    # Zebra Stripe Alternating Data Processing Rows
    for r_idx, dataset in enumerate(records):
        for c_idx, text_val in enumerate(dataset):
            cell = grid.cell(r_idx + 1, c_idx)
            cell.text = text_val
            cell.fill.solid()
            # Clean light/white alternating block lines
            stripe_color = RGBColor(241, 245, 249) if r_idx % 2 == 0 else COLOR_WHITE
            cell.fill.fore_color.rgb = stripe_color
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_FAMILY
                p.font.size = Pt(10)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PRIMARY_DARK
                else:
                    p.font.color.rgb = COLOR_TEXT_MAIN
                p.alignment = PP_ALIGN.LEFT if c_idx == 2 else PP_ALIGN.CENTER


# ════════════════════ INJECTING SLIDE DECKS GENERATION ════════════════════

# --- SLIDE 1: LANDING DISPLAY COVER (Dark) ---
s1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s1, COLOR_PRIMARY_DARK)
draw_header_footer(s1, None, 1, is_dark_layout=True)

cover_context = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(0.4))
p_cc = cover_context.text_frame.paragraphs[0]
p_cc.text = "PROJECT REVIEW PRESENTATION"
p_cc.alignment = PP_ALIGN.CENTER
p_cc.font.name = FONT_FAMILY
p_cc.font.size = Pt(11)
p_cc.font.bold = True
p_cc.font.color.rgb = COLOR_ACCENT

cover_h1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.33), Inches(2.0))
tf_ch1 = cover_h1.text_frame
tf_ch1.word_wrap = True
p_ch1 = tf_ch1.paragraphs[0]
p_ch1.text = "ONLINE HOUSE RENTAL SYSTEM"
p_ch1.alignment = PP_ALIGN.CENTER
p_ch1.font.name = FONT_FAMILY
p_ch1.font.size = Pt(44)
p_ch1.font.bold = True
p_ch1.font.color.rgb = COLOR_WHITE
p_ch1.space_after = Pt(12)

p_ch2 = tf_ch1.add_paragraph()
p_ch2.text = '"Seamless Property Search, Booking & Management Platform"'
p_ch2.alignment = PP_ALIGN.CENTER
p_ch2.font.name = FONT_FAMILY
p_ch2.font.size = Pt(15)
p_ch2.font.color.rgb = RGBColor(203, 213, 225)
p_ch2.font.italic = True

middle_line = s1.shapes.add_shape(1, Inches(5.66), Inches(4.3), Inches(2.0), Inches(0.04))
middle_line.fill.solid()
middle_line.fill.fore_color.rgb = COLOR_ACCENT
middle_line.line.fill.background()

tech_stripe = s1.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.33), Inches(1.0))
tf_ts = tech_stripe.text_frame
p_ts = tf_ts.paragraphs[0]
p_ts.text = "Backend: PHP Core & Sessions    |    Database: MySQL Relational    |    Frontend: Bootstrap 5 & JS"
p_ts.alignment = PP_ALIGN.CENTER
p_ts.font.name = FONT_FAMILY
p_ts.font.size = Pt(12.5)
p_ts.font.bold = True
p_ts.font.color.rgb = COLOR_WHITE


# --- SLIDE 2: BUSINESS OVERVIEW VISIONS (Light) ---
s2 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s2, COLOR_LIGHT_BG)
draw_header_footer(s2, "Project Overview", 2)

overview_bullet_points = [
    "Direct-to-Consumer Ecosystem: Connects property owners/landlords directly with prospective tenants, bypassing standard middlemen and platform commissions.",
    "Dynamic Listings: Provides unified, real-time catalog storage categorized across multiple structural classes: Apartments, Houses, Villas, and Studios.",
    "Integrated Booking Pipelines: Facilitates end-to-end booking inquiries mapping strict target dates (Check-In & Check-Out validation constraints).",
    "Multi-Dimensional Queries: Enables advanced filtering protocols linking pricing thresholds, geographical areas, and capacity counts."
]
for i, statement in enumerate(overview_bullet_points):
    draw_bullet_point(s2, i + 1, statement, 1.8 + (i * 1.2))


# --- SLIDE 3: SYSTEMS INFRASTRUCTURE LAYOUT (Light Grid) ---
s3 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s3, COLOR_LIGHT_BG)
draw_header_footer(s3, "Technology Architecture", 3)

draw_container_card(s3, 0.6, 1.6, 3.7, 4.8, "Frontend Tier", [
    "HTML5 & Custom CSS3 Elements",
    "Bootstrap 5 Responsive Grid System",
    "JavaScript Interface Component Actions",
    "Integrated FontAwesome Vector Assets"
])
draw_container_card(s3, 4.81, 1.6, 3.7, 4.8, "Backend Tier", [
    "PHP Core Engine Runtime",
    "Strict Server-Side Input Validation",
    "File-Upload Handling & Sanitization Module",
    "Native Session State Statekeeping Management"
])
draw_container_card(s3, 9.03, 1.6, 3.7, 4.8, "Database Tier", [
    "MySQL Relational Database Server",
    "7 Fully Normalized Core Data Entities",
    "Primary-to-Foreign Key Configuration",
    "Enforced Cascading Referential Integrity"
])


# --- SLIDE 4: CENTRALIZED CONTROL MODULES (Asymmetric) ---
s4 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s4, COLOR_LIGHT_BG)
admin_data = [
    {"title": "User Supervision", "desc": "The centralized administrator holds global rights to audit landlord and tenant operations. They oversee registration pipelines, review system activation, and handle credential updates or account state blocks."},
    {"title": "Property Audit & Verify", "desc": "Every landlord-onboarded housing asset routes to a pending verification queue. Admins evaluate documentation, confirm structural capacities, and green-light properties to make them public."}
]
draw_asymmetric_sidebar(s4, COLOR_PRIMARY_DARK, "Admin\nModule\nFeatures", admin_data, 4)


# --- SLIDE 5: RESIDENT LANDLORD MANAGEMENT (Asymmetric) ---
s5 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s5, COLOR_LIGHT_BG)
owner_data = [
    {"title": "Property Onboarding", "desc": "Enables property owners to list rental inventories. Landlords specify housing titles, geographical locations, monthly rental pricing, room configurations, descriptive features, and upload structural photographs."},
    {"title": "Reservation Monitoring", "desc": "Owners monitor lease bookings for their respective listings. They can view tenant-submitted check-in and check-out timelines, audit tenant information, and approve/reject bookings dynamically."}
]
draw_asymmetric_sidebar(s5, COLOR_TEAL_PANEL, "Owner\nModule\nFeatures", owner_data, 5)


# --- SLIDE 6: TENANT DISCOVERY OPERATIONS (Asymmetric) ---
s6 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s6, COLOR_LIGHT_BG)
tenant_data = [
    {"title": "Listing Discovery", "desc": "Allows tenants to explore live listings. Clients can sort available accommodations dynamically by pricing ranges, capacity configurations, geographical locations, and layout types."},
    {"title": "Stay Tracking & History", "desc": "Tenants trigger custom booking requests by specifying structural check-in and check-out dates. They can monitor reservation states (Approved/Pending/Rejected) and audit stay histories."}
]
draw_asymmetric_sidebar(s6, COLOR_AMBER_PANEL, "Tenant\nModule\nFeatures", tenant_data, 6)


# --- SLIDE 7: INTERCONNECTED MODELLING DIVIDER (Dark Section) ---
s7 = prs.slides.add_slide(prs.slide_layouts[6])
build_section_divider(s7, "UML Diagrams", "System Architecture & Dynamic Model Mappings", 7)


# --- SLIDE 8: DATA RELATIONAL MAP DIVIDER (Dark Section) ---
s8 = prs.slides.add_slide(prs.slide_layouts[6])
build_section_divider(s8, "Database Schema", "Relational Mapping of the 'house_rental' Database", 8)


# Standard Columns Structuring
schema_headers = ["Column File Header", "Data Structural Type", "Functional Data Field Description"]

# --- SLIDE 9: ADMIN & SECURITY ACCESS TABLES ---
s9 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s9, COLOR_LIGHT_BG)
draw_header_footer(s9, "Schema: Admin & User Tables", 9)

admin_records = [
    ["id", "INT (PK)", "Unique administrative profile sequential primary tracking key"],
    ["username", "VARCHAR", "Alpha-numeric login account configuration profile string"],
    ["password", "VARCHAR", "Secure server hash-stored key system credential validation"]
]
build_relational_table(s9, 0.6, 1.4, 12.13, 2.0, "Table: admin", schema_headers, admin_records)

user_records = [
    ["id", "INT (PK)", "Unique registered tenant profile identity sequence record code"],
    ["username", "VARCHAR", "Public display identifier for tenant operations dashboard interface"],
    ["password", "VARCHAR", "Encrypted string password hash for profile security shielding"],
    ["doc", "VARCHAR", "Server-relative asset address path directing to tenant uploaded identity logs"],
    ["email", "VARCHAR", "Primary digital contact point node for system messaging feeds"]
]
build_relational_table(s9, 0.6, 4.0, 12.13, 2.6, "Table: user", schema_headers, user_records)


# --- SLIDE 10: PROPERTY OWNERSHIP LOGISTICS TABLES ---
s10 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s10, COLOR_LIGHT_BG)
draw_header_footer(s10, "Schema: Owner & Category Tables", 10)

owner_records = [
    ["id", "INT (PK)", "Unique registered property investor onboarding signature database code"],
    ["name", "VARCHAR", "Legal corporate or personal representative identity title string"],
    ["email", "VARCHAR", "Owner designated electronic communications link validation gateway"],
    ["doc", "VARCHAR", "Local file storage address pointing to corporate license proof tracking metadata"],
    ["status", "TINYINT", "Binary workflow toggle code tracking active verification (1=Active, 0=Pending)"]
]
build_relational_table(s10, 0.6, 1.4, 12.13, 2.6, "Table: owner", schema_headers, owner_records)

cat_records = [
    ["id", "INT (PK)", "Structural categorization primary system processing directory tracking key"],
    ["name", "VARCHAR", "Property structural design framework type string (Apartment, Villa, House, Studio)"]
]
build_relational_table(s10, 0.6, 4.6, 12.13, 2.0, "Table: category", schema_headers, cat_records)


# --- SLIDE 11: HOUSES REGISTRY CATALOGUE (Full Space Grid) ---
s11 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s11, COLOR_LIGHT_BG)
draw_header_footer(s11, "Schema: Houses Table Structure", 11)

house_records = [
    ["id", "INT (PK)", "Auto-incrementing asset identity key index mapping the system inventory ledger"],
    ["owner_id", "INT (FK)", "Links tracking node backward to owner.id identifying structural manager entity"],
    ["title", "VARCHAR", "Marketing display heading and interface text label for target property"],
    ["location", "VARCHAR", "Geographical mapping data details specifying street physical location data"],
    ["price", "DECIMAL(12,2)", "Quantified monetary cost structure detailing expected system flat rate monthly rental price"],
    ["amenities", "VARCHAR", "Aggregated character line text logging utility arrays like WiFi, Parking, AC or Gym modules"]
]
build_relational_table(s11, 0.6, 1.4, 12.13, 5.0, "Table: houses", schema_headers, house_records)


# --- SLIDE 12: FINANCIAL BOOKING LEDGER TRANSACTION TABLES ---
s12 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s12, COLOR_LIGHT_BG)
draw_header_footer(s12, "Schema: Bookings Table Structure", 12)

booking_records = [
    ["id", "INT (PK)", "Unique system confirmation receipt ticket code identifier tracking the booking request"],
    ["user_id", "INT (FK)", "Connects specific system record entry downward to requesting consumer target via user.id"],
    ["house_id", "INT (FK)", "Maps active system interface lease parameters directly onto inventory via houses.id link"],
    ["checkin_date", "DATE", "Designated start system date constraint verifying initial active calendar date entry limit"],
    ["checkout_date", "DATE", "Designated end system date constraint verifying final target contract timeline exit checkpoint"],
    ["status", "TINYINT", "Workflow lifecycle pipeline processing state code flags tracking Approval state boundaries"]
]
build_relational_table(s12, 0.6, 1.4, 12.13, 5.0, "Table: bookings", schema_headers, booking_records)


# --- SLIDE 13: SANDBOX ENTRY AUDITING BUFFERS (Light) ---
s13 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s13, COLOR_LIGHT_BG)
draw_header_footer(s13, "Schema: Property Buffer Table", 13)

prop_records = [
    ["id", "INT (PK)", "Unique pipeline transactional processing index key routing new application files"],
    ["owner_id", "INT (FK)", "Identifies requesting provider agent submitting request relative to owner.id table"],
    ["title", "VARCHAR", "Proposed property listing name label submitted for evaluation review"],
    ["price", "DECIMAL(12,2)", "Proposed monthly financial retail contract baseline price structure"],
    ["status", "TINYINT", "Validation checkpoint indicator registering administrative active approval gate results"]
]
build_relational_table(s13, 0.6, 1.4, 12.13, 4.2, "Table: property", schema_headers, prop_records)

note_box = s13.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.4))
tf_nb = note_box.text_frame
tf_nb.word_wrap = True
p_nb = tf_nb.paragraphs[0]
p_nb.text = "Note: The 'property' entity functions exclusively as an isolated system staging sandbox buffer. Data fields are held for staging evaluation until global administrative clearance actions migrate standard parameters into the live operational 'houses' catalogue database ledger grid."
p_nb.font.name = FONT_FAMILY
p_nb.font.size = Pt(9.5)
p_nb.font.italic = True
p_nb.font.color.rgb = COLOR_TEXT_MUTED


# --- SLIDE 14: THANK YOU WRAP DISPLAY (Dark Closing) ---
s14 = prs.slides.add_slide(prs.slide_layouts[6])
apply_solid_background(s14, COLOR_PRIMARY_DARK)
draw_header_footer(s14, None, 14, is_dark_layout=True)

end_box = s14.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
tf_eb = end_box.text_frame
tf_eb.word_wrap = True

p_e1 = tf_eb.paragraphs[0]
p_e1.text = "THANK YOU!"
p_e1.alignment = PP_ALIGN.CENTER
p_e1.font.name = FONT_FAMILY
p_e1.font.size = Pt(44)
p_e1.font.bold = True
p_e1.font.color.rgb = COLOR_WHITE
p_e1.space_after = Pt(12)

p_e2 = tf_eb.add_paragraph()
p_e2.text = "Online House Rental System Platform — Connecting People with Homes"
p_e2.alignment = PP_ALIGN.CENTER
p_e2.font.name = FONT_FAMILY
p_e2.font.size = Pt(14)
p_e2.font.color.rgb = RGBColor(148, 163, 184)

end_line = s14.shapes.add_shape(1, Inches(5.66), Inches(4.3), Inches(2.0), Inches(0.04))
end_line.fill.solid()
end_line.fill.fore_color.rgb = COLOR_ACCENT
end_line.line.fill.background()

sub_end_box = s14.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.33), Inches(1.0))
p_seb = sub_end_box.text_frame.paragraphs[0]
p_seb.text = "Questions, Feedback, or System Suggestions?"
p_seb.alignment = PP_ALIGN.CENTER
p_seb.font.name = FONT_FAMILY
p_seb.font.size = Pt(13)
p_seb.font.bold = True
p_seb.font.color.rgb = COLOR_ACCENT


# --- PROCESS EXPORT FILENAME ---
output_file = "Online_House_Rental_Beautiful_Presentation.pptx"
prs.save(output_file)
print(f"Design-Optimized Presentation saved as: {os.path.abspath(output_file)}")
